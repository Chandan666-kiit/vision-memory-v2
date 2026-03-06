"""
Generate a colourful PowerPoint for the Vision + Memory project with
LangGraph details and screenshot placeholders. Run: python create_presentation.py
Requires: pip install python-pptx pillow
"""
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------- Colour palette ----------
DARK_BLUE = RGBColor(0x1E, 0x3A, 0x5F)       # title / summary
TEAL = RGBColor(0x0D, 0x94, 0x88)             # overview / used
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF7, 0xF7)        # light teal-grey
GREEN_BG = RGBColor(0xE6, 0xF7, 0xF2)        # used feature
ORANGE_BG = RGBColor(0xFF, 0xF4, 0xE6)       # not used
ACCENT_GREEN = RGBColor(0x00, 0x7A, 0x3D)
ACCENT_ORANGE = RGBColor(0xE6, 0x51, 0x00)
ACCENT_PURPLE = RGBColor(0x5E, 0x35, 0xB1)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0x88, 0x88, 0x88)

ASSETS_DIR = Path(__file__).parent / "ppt_assets"
ASSETS_DIR.mkdir(exist_ok=True)


def set_slide_background(slide, color, width_inches=10, height_inches=7.5):
    """Set slide background to a solid colour (full-slide rectangle so it's always visible)."""
    rect = slide.shapes.add_shape(
        1, 0, 0, Inches(width_inches), Inches(height_inches)
    )  # 1 = rectangle
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    for _ in range(15):
        try:
            rect.z_order.send_backward()
        except Exception:
            break


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BLUE)
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.4))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1.2))
        p2 = tx2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(26)
        p2.font.color.rgb = RGBColor(0xAD, 0xC5, 0xE0)
        p2.alignment = PP_ALIGN.CENTER


def add_section_header(prs, title, subtitle=""):
    """Coloured section header slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, TEAL)
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.2))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
        p2 = tx2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0xE0, 0xF7, 0xF5)
        p2.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, bullets, subbullets=None, bg_color=LIGHT_BG, title_color=DARK_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg_color)
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.85))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = title_color
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(9), Inches(5.6))
    tf = body.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = item
        para.font.size = Pt(17)
        para.font.color.rgb = GRAY
        para.space_after = Pt(6)
        if subbullets and i < len(subbullets) and subbullets[i]:
            for sub in subbullets[i]:
                sp = tf.add_paragraph()
                sp.text = "  ▪ " + sub
                sp.font.size = Pt(14)
                sp.font.color.rgb = LIGHT_GRAY
                sp.space_after = Pt(2)
    return slide


def add_feature_slide(prs, feature_name, used, description, details, extra=None):
    """Detailed feature slide with green (used) or orange (not used) accent."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = GREEN_BG if used else ORANGE_BG
    set_slide_background(slide, bg)
    accent = ACCENT_GREEN if used else ACCENT_ORANGE
    # Title
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.75))
    p = tx.text_frame.paragraphs[0]
    p.text = feature_name
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = accent
    # Badge
    status = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(9), Inches(0.45))
    sp = status.text_frame.paragraphs[0]
    sp.text = "✓ USED IN THIS PROJECT" if used else "○ NOT USED (possible extension)"
    sp.font.size = Pt(13)
    sp.font.italic = True
    sp.font.color.rgb = accent
    # Description
    desc = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(1.5))
    dp = desc.text_frame.paragraphs[0]
    dp.text = description
    dp.font.size = Pt(16)
    dp.font.color.rgb = GRAY
    dp.space_after = Pt(10)
    # Details
    top = 3.2 if not extra else 3.0
    det = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(3.8))
    dtf = det.text_frame
    for i, line in enumerate(details):
        para = dtf.paragraphs[0] if i == 0 else dtf.add_paragraph()
        para.text = "• " + line
        para.font.size = Pt(14)
        para.font.color.rgb = GRAY
        para.space_after = Pt(3)
    if extra:
        ep = dtf.add_paragraph()
        ep.text = extra
        ep.font.size = Pt(12)
        ep.font.italic = True
        ep.font.color.rgb = LIGHT_GRAY
    return slide


def add_screenshot_slide(prs, title, img_path, caption):
    """Slide with a large image (screenshot placeholder or real)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, RGBColor(0xEE, 0xEE, 0xEE))
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    if img_path and os.path.isfile(img_path):
        slide.shapes.add_picture(str(img_path), Inches(0.8), Inches(1.1), width=Inches(8.4))
    else:
        # Placeholder box
        ph = slide.shapes.add_shape(1, Inches(1.5), Inches(1.5), Inches(7), Inches(4))  # rectangle
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        ph.line.color.rgb = RGBColor(0x99, 0x99, 0x99)
        tx2 = slide.shapes.add_textbox(Inches(2), Inches(3.2), Inches(6), Inches(1.2))
        tp = tx2.text_frame.paragraphs[0]
        tp.text = "Replace with your screenshot"
        tp.font.size = Pt(20)
        tp.font.color.rgb = LIGHT_GRAY
        tp.alignment = PP_ALIGN.CENTER
    cap = slide.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(9), Inches(0.6))
    cp = cap.text_frame.paragraphs[0]
    cp.text = caption
    cp.font.size = Pt(14)
    cp.font.color.rgb = GRAY
    cp.alignment = PP_ALIGN.CENTER


def create_placeholder_screenshot_images():
    """Create simple placeholder images for screenshots (replace with real ones later)."""
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        return []
    W, H = 800, 450
    fonts = []
    try:
        fonts.append(ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 28))
        fonts.append(ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 18))
    except Exception:
        fonts = [ImageFont.load_default(), ImageFont.load_default()]
    out_paths = []
    # 1. Upload screen
    img1 = Image.new("RGB", (W, H), (0xE8, 0xF4, 0xF4))
    d = ImageDraw.Draw(img1)
    d.rectangle([20, 20, W - 20, 80], fill=(0x0D, 0x94, 0x88), outline=(0x0A, 0x6B, 0x60))
    d.text((W // 2 - 120, 42), "Vision & Memory - Streamlit", fill=(255, 255, 255), font=fonts[0])
    d.rectangle([80, 120, W - 80, 200], fill=(255, 255, 255), outline=(0x0D, 0x94, 0x88))
    d.text((W // 2 - 180, 155), "📤 Upload an image (person face)", fill=(0x33, 0x33, 0x33), font=fonts[1])
    d.text((W // 2 - 100, 320), "Screenshot 1: File upload screen", fill=(0x66, 0x66, 0x66), font=fonts[1])
    p1 = ASSETS_DIR / "screenshot1_upload.png"
    img1.save(p1)
    out_paths.append(p1)
    # 2. Analysis (parallel: 4 metrics) + name form
    img2 = Image.new("RGB", (W, H), (0xF5, 0xFA, 0xFA))
    d2 = ImageDraw.Draw(img2)
    d2.rectangle([30, 25, W - 30, 185], fill=(255, 255, 255), outline=(0x0D, 0x94, 0x88))
    d2.text((50, 45), "Analysis (parallel: person, glasses, emotion, age)", fill=(0x1E, 0x3A, 0x5F), font=fonts[0])
    d2.text((50, 85), "Person: Yes   Glasses: Yes   Emotion: happy   Age: 25-30", fill=(0x00, 0x7A, 0x3D), font=fonts[1])
    d2.rectangle([30, 200, 770, 320], fill=(0xFF, 0xF9, 0xE6), outline=(0xE6, 0x51, 0x00))
    d2.text((50, 220), "We don't recognize you yet. Enter your name below.", fill=(0x66, 0x66, 0x66), font=fonts[1])
    d2.rectangle([50, 255, 400, 295], fill=(255, 255, 255), outline=(0xCC, 0xCC, 0xCC))
    d2.text((65, 265), "Your name", fill=(0x99, 0x99, 0x99), font=fonts[1])
    d2.rectangle([420, 260, 550, 295], fill=(0x0D, 0x94, 0x88))
    d2.text((455, 268), "Submit name", fill=(255, 255, 255), font=fonts[1])
    d2.text((W // 2 - 160, 400), "Screenshot 2: Analysis & name form", fill=(0x66, 0x66, 0x66), font=fonts[1])
    p2 = ASSETS_DIR / "screenshot2_analysis_form.png"
    img2.save(p2)
    out_paths.append(p2)
    # 3. Greeting
    img3 = Image.new("RGB", (W, H), (0xE8, 0xF7, 0xF2))
    d3 = ImageDraw.Draw(img3)
    d3.rectangle([80, 120, W - 80, 280], fill=(255, 255, 255), outline=(0x00, 0x7A, 0x3D))
    d3.text((W // 2 - 180, 160), "Greeting", fill=(0x1E, 0x3A, 0x5F), font=fonts[0])
    d3.text((60, 210), "Good morning Alex! — 👓 Looking sharp with specs,", fill=(0x00, 0x7A, 0x3D), font=fonts[1])
    d3.text((60, 235), "you look happy, around 25-30.", fill=(0x00, 0x7A, 0x3D), font=fonts[1])
    d3.text((W // 2 - 100, 360), "Screenshot 3: Greeting output", fill=(0x66, 0x66, 0x66), font=fonts[1])
    p3 = ASSETS_DIR / "screenshot3_greeting.png"
    img3.save(p3)
    out_paths.append(p3)
    return out_paths


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create placeholder screenshots (user can replace with real PNGs in ppt_assets/)
    screenshots = create_placeholder_screenshot_images()
    s1 = screenshots[0] if len(screenshots) > 0 else None
    s2 = screenshots[1] if len(screenshots) > 1 else None
    s3 = screenshots[2] if len(screenshots) > 2 else None

    # ---- 1. Title ----
    add_title_slide(
        prs,
        "Vision & Memory Project",
        "LangGraph features, architecture & app output"
    )

    # ---- 2. Project overview (detailed) ----
    add_section_header(prs, "Project overview", "What the application does")
    add_content_slide(
        prs,
        "What is this project?",
        [
            "A web app where you upload a photo; the system runs four vision analyses in parallel (person, glasses, emotion, age), merges results, then either recognizes the person from memory or asks for their name and remembers it.",
            "Flow: Upload image → four parallel nodes (detect_person, detect_glasses, detect_emotion, detect_age) → merge_vision → identify (FAISS or need_name) → greet (personalised message with emotion/age).",
            "Frontend: Streamlit (upload, 4 metrics, name form when needed, greeting). Backend: LangGraph StateGraph with parallel nodes + merge, and custom FAISS long-term memory.",
        ],
        [
            ["Fan-out from START to 4 nodes; fan-in to merge_vision; then identify → greet → END"],
            ["State: image_path, person, glasses, emotion, age_estimate, person_name, need_name, greeting_message"],
            ["Run: streamlit run streamlit_app.py"],
        ],
        bg_color=LIGHT_BG,
        title_color=DARK_BLUE,
    )

    # ---- 3. Feature descriptions (what each part does) ----
    add_section_header(prs, "Feature descriptions", "What each component does")
    add_content_slide(
        prs,
        "Feature descriptions — Vision & merge",
        [
            "Person detection — Parallel node. Asks GPT-4o-mini: 'Is there a human face?' Returns true/false. Decides whether we show a greeting or skip.",
            "Glasses detection — Parallel node. 'Does the person wear glasses?' Returns true/false. Shown in UI and in greeting (e.g. 'Looking sharp with specs').",
            "Emotion detection — Parallel node. 'What is the dominant emotion?' Returns e.g. happy, neutral. Shown in analysis and in the greeting.",
            "Age estimation — Parallel node. 'Estimate age range.' Returns e.g. 25-30, 30-40. Shown in analysis and in greeting ('around 25-30').",
            "Merge (merge_vision) — Runs after all four parallel nodes. Combines person, glasses, emotion, age_estimate into response_text (JSON) for the rest of the pipeline and UI. All four vision nodes run in parallel from START; merge runs once after they complete.",
        ],
        None,
        bg_color=LIGHT_BG,
        title_color=DARK_BLUE,
    )
    add_content_slide(
        prs,
        "Feature descriptions — Identify, greet & memory",
        [
            "Identify — If no person, returns. If person_name already in state (from form), adds 'Name: X | Glasses: Y' to FAISS and continues. Otherwise sets need_name=True so Streamlit shows the name form.",
            "Greet — Builds greeting_message from person_name, glasses, emotion, age_estimate. Example: 'Good morning Alex! — Looking sharp with specs, you look happy, around 25-30.'",
            "FAISS long-term memory — Custom vector store (not LangGraph). Stores lines like 'Name: Alex | Glasses: True'. Retrieved by similarity in identify. Persisted to faiss_index/ so names survive restarts.",
            "Store (InMemoryStore) — LangGraph built-in. Passed as graph.compile(store=store). Used for graph-level storage; ready for future tool/agent state. identify_person uses LongTermMemory.retrieve_memory() and add_memory().",
        ],
        None,
        bg_color=LIGHT_BG,
        title_color=DARK_BLUE,
    )

    # ---- 4. App output screenshots ----
    add_section_header(prs, "App output", "Screenshots from the Streamlit UI")
    add_screenshot_slide(
        prs,
        "Screenshot 1: File upload",
        s1,
        "User uploads an image; the app saves it and runs the LangGraph pipeline.",
    )
    add_screenshot_slide(
        prs,
        "Screenshot 2: Analysis & name form",
        s2,
        "After parallel analysis (person, glasses, emotion, age), if new, the app shows all four metrics and asks for name.",
    )
    add_screenshot_slide(
        prs,
        "Screenshot 3: Greeting",
        s3,
        "After name submit: added to FAISS; greeting can include glasses, emotion, and age estimate.",
    )

    # ---- 5. LangGraph features at a glance ----
    add_section_header(prs, "LangGraph features", "What is used in this project")
    add_content_slide(
        prs,
        "LangGraph features in this project",
        [
            "State graph & typed state — StateGraph(VisionState) with image_path, person, glasses, emotion, age_estimate, person_name, response_text, need_name, greeting_message. All nodes read/update this state.",
            "Parallel nodes (fan-out / fan-in) — Four nodes run in parallel from START: detect_person, detect_glasses, detect_emotion, detect_age. All feed into merge_vision, then identify → greet → END. Shows advanced orchestration.",
            "Store (InMemoryStore) — graph.compile(store=store). LangGraph Store API for optional graph-level storage (e.g. future tool/agent state).",
            "Long-term memory (FAISS) — Custom FAISS + OpenAI embeddings inside the identify node; stores/retrieves name + glasses. Persisted to faiss_index/.",
        ],
        [
            ["StateGraph(VisionState), get_initial_state(), app.invoke(state)"],
            ["add_edge(START, 'detect_*') × 4; add_edge('detect_*', 'merge_vision') × 4; then merge → identify → greet → END"],
            ["InMemoryStore; compile(store=store)"],
            ["LongTermMemory in identify_person"],
        ],
        bg_color=GREEN_BG,
        title_color=TEAL,
    )

    # ---- 6. State & StateGraph (detailed) ----
    add_content_slide(
        prs,
        "State graph & typed state — in detail",
        [
            "What it is: A StateGraph shares one state object across nodes. Each node receives and returns (partial) state; LangGraph merges updates so downstream nodes see the latest values.",
            "What we use: VisionState includes image_path, person, glasses, emotion, age_estimate, person_name, response_text, need_name, greeting_message. The four parallel nodes each write only their keys (person, glasses, emotion, age_estimate); merge_vision sets response_text; identify and greet use the rest. get_initial_state() and app.invoke(initial_state) drive the run.",
            "Why it matters: Partial updates from parallel nodes are merged; merge_vision then runs with the combined state. The frontend reads need_name, greeting_message, emotion, age_estimate from the final state.",
        ],
        [
            ["graph = StateGraph(VisionState)"],
            ["Parallel nodes return e.g. {\"person\": True} or {\"emotion\": \"happy\"}; merge merges all"],
            ["Streamlit: result['need_name'], result['greeting_message'], result['emotion'], result['age_estimate']"],
        ],
        bg_color=GREEN_BG,
        title_color=ACCENT_GREEN,
    )

    # ---- 7. Parallel nodes & merge (detailed) ----
    add_content_slide(
        prs,
        "Parallel nodes & merge — in detail",
        [
            "What they are: LangGraph can fan out from START (or any node) to multiple nodes; those nodes run in parallel in one superstep. A merge node has multiple incoming edges; it runs after all parallel nodes complete, with the merged state.",
            "What we use: (1) Four parallel nodes: detect_person, detect_glasses, detect_emotion, detect_age — each encodes the image and calls GPT-4o-mini with a different prompt, returns only its key(s). (2) merge_vision — builds response_text from person, glasses, emotion, age_estimate. (3) identify — FAISS or need_name. (4) greet — greeting_message using name, glasses, emotion, age.",
            "Edges: add_edge(START, 'detect_person'), add_edge(START, 'detect_glasses'), add_edge(START, 'detect_emotion'), add_edge(START, 'detect_age'); then all four → merge_vision → identify → greet → END.",
        ],
        [
            ["Fan-out: START → 4 nodes in parallel; fan-in: 4 nodes → merge_vision"],
            ["Each detect_* node returns e.g. {\"person\": True} or {\"emotion\": \"happy\"}; state is merged"],
            ["Shows advanced orchestration: parallel execution then merge"],
        ],
        bg_color=GREEN_BG,
        title_color=ACCENT_GREEN,
    )

    # ---- 8. Store (detailed) ----
    add_feature_slide(
        prs,
        "Store (InMemoryStore)",
        True,
        "LangGraph’s Store API provides key-value storage attached to the graph. InMemoryStore is a simple in-memory implementation. We pass it when compiling: app = graph.compile(store=store). It allows the graph to use a store for namespaced data (e.g. for tools or multi-turn agent state). In this project we do not read/write the store in our nodes; we use it so the compiled app is ready for future extensions (e.g. storing conversation or tool results).",
        [
            "from langgraph.store.memory import InMemoryStore",
            "store = InMemoryStore(); app = graph.compile(store=store)",
            "Store is separate from our FAISS long-term memory (which is custom and used inside identify_person).",
        ],
        extra="Future use: persist tool outputs or agent context in the store.",
    )

    # ---- 9. Long-term memory FAISS (detailed) ----
    add_content_slide(
        prs,
        "Long-term memory (FAISS) — what it is and how we use it",
        [
            "What it is: A custom component, not a LangGraph built-in. We use FAISS (via LangChain) as a vector store: text is embedded with OpenAI embeddings and stored; we retrieve by similarity (e.g. 'Name:' to get stored name+glasses lines).",
            "What we use: LongTermMemory class with add_memory(text) and retrieve_memory(query, k). We store lines like 'Name: Alex | Glasses: True'. The identify node calls retrieve_memory('Name:', k=10), parses names and glasses, and either matches (we removed weak glasses-only matching) or sets need_name=True so the UI asks; when the user submits a name we re-invoke with person_name in state and identify_person then calls add_memory.",
            "Persistence: The index is saved under faiss_index/ so remembered names survive restarts.",
        ],
        [
            ["FAISS + OpenAIEmbeddings; Document(page_content=text)"],
            ["Used only inside the identify_person node"],
        ],
        bg_color=GREEN_BG,
        title_color=TEAL,
    )

    # ---- 10. Parallelization (USED) ----
    add_feature_slide(
        prs,
        "Parallel nodes (fan-out / fan-in)",
        True,
        "LangGraph supports running multiple nodes in parallel: a node (or START) can have multiple outgoing edges; all target nodes run in the same superstep. We use this for vision: four nodes (person, glasses, emotion, age) run in parallel, then a merge node combines their outputs into state. This shows advanced orchestration — parallel execution then merge.",
        [
            "add_edge(START, 'detect_person'), add_edge(START, 'detect_glasses'), add_edge(START, 'detect_emotion'), add_edge(START, 'detect_age')",
            "All four → merge_vision; merge_vision builds response_text from the four results",
            "Each parallel node returns only its keys; LangGraph merges partial state updates",
        ],
        extra="Same pattern can extend to more parallel analyses (e.g. attire, background).",
    )

    # ---- 11. Agent / tools (not used) ----
    add_feature_slide(
        prs,
        "Agent / tools & human-in-the-loop",
        False,
        "LangGraph can power agents that call tools (e.g. bind_tools with an LLM and tool-execution nodes). We do not use an agent loop or tool calls. Human-in-the-loop here means: when we need a name, we set need_name=True and return; the Streamlit UI shows a form, the user submits; we then invoke the graph again with person_name in the initial state. So the 'human' step is outside the graph (in the frontend), not inside a LangGraph node.",
        [
            "No bind_tools or tool-execution nodes",
            "Human input: name collected in Streamlit; second invoke(state) with person_name set",
        ],
        extra="To add agents: introduce a tool-calling node and connect to tools.",
    )

    # ---- 12. Summary ----
    add_section_header(prs, "Summary", "What is used vs not used")
    add_content_slide(
        prs,
        "Summary",
        [
            "Used: StateGraph(VisionState) with emotion & age_estimate; parallel nodes (detect_person, detect_glasses, detect_emotion, detect_age) with fan-out from START and fan-in to merge_vision; then identify → greet → END. compile(store=InMemoryStore), invoke(initial_state). FAISS long-term memory in the identify node. Greeting uses person, glasses, emotion, age.",
            "Not used: Conditional edges, agent/tool nodes, checkpointer. Human-in-the-loop is in the UI (Streamlit form + re-invoke), not inside the graph.",
            "Screenshots: Replace ppt_assets/screenshot*.png with your own and re-run create_presentation.py.",
        ],
        [
            ["State, parallel nodes, merge, identify, greet, store, invoke — all used"],
            ["Conditional edges, agent, checkpointer — not used"],
            ["Replace ppt_assets/*.png with real screenshots and re-run script"],
        ],
        bg_color=RGBColor(0xE8, 0xED, 0xF5),
        title_color=DARK_BLUE,
    )

    out_path = "Vision_Memory_LangGraph_Features.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")
    if screenshots:
        print(f"Placeholder screenshots saved in: {ASSETS_DIR}")
    else:
        print("Install Pillow (pip install pillow) to generate placeholder screenshots.")


if __name__ == "__main__":
    main()
